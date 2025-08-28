# app_streamlit_atenciones_v16.py
# Dashboard de Pacientes – organizado en secciones y funciones

import streamlit as st
import pandas as pd
import numpy as np
import altair as alt
from io import BytesIO
from datetime import datetime

# ===== PowerPoint =====
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor

# ===== NUEVO: para lectura por URL =====
import requests

# ===================== CONFIGURACIÓN GLOBAL =====================
st.set_page_config(page_title="Dashboard de Pacientes", page_icon="🩺", layout="wide")

# Tema oscuro legible
st.markdown("""
<style>
  h1, h2, h3, label, .stMarkdown, .stCaption { color:#F8FAFC !important; }
  h1 { font-weight:800; letter-spacing:-.3px; font-family:ui-sans-serif,system-ui,"Segoe UI",Roboto; }
  .section-title { font-size:1.05rem; font-weight:700; color:#E5E7EB; margin:6px 0 10px; }
  .kpi-card{
    padding:14px 16px; border-radius:14px;
    background:#0b1220; border:1px solid #1f2937; 
    box-shadow:0 2px 12px rgba(0,0,0,.35);
  }
  .kpi-title{ font-size:.9rem; color:#CBD5E1; margin-bottom:6px; }
  .kpi-value{ font-size:1.7rem; font-weight:800; color:#F8FAFC; }
</style>
""", unsafe_allow_html=True)

PALETTE = ["#60A5FA","#34D399","#FBBF24","#A78BFA","#F87171","#22D3EE",
           "#F472B6","#4ADE80","#93C5FD","#F59E0B"]

REQUIRED_COLS = [
    "NUMERO_PACIENTE","FECHA_ATENCION","GENERO","EDAD","SEDE",
    "EPS","PROGRAMA","ESPECIALIDAD","PRESTADOR","ENFERMERA","MES_ATENCION","CIUDAD"
]
DATE_COLS = ["FECHA_ATENCION"]


# ===================== UTILIDADES =====================
def normalize_cols(df: pd.DataFrame) -> pd.DataFrame:
    return df.rename(columns={c: c.strip() for c in df.columns})

def ensure_columns(df: pd.DataFrame) -> pd.DataFrame:
    for c in REQUIRED_COLS:
        if c not in df.columns:
            df[c] = pd.NA
    return df

def excel_serial_to_datetime(series: pd.Series):
    ser_num = pd.to_numeric(series, errors="coerce")
    base = pd.to_datetime("1899-12-30")
    return base + pd.to_timedelta(ser_num, unit="D")

def robust_parse_datetime(s: pd.Series) -> pd.Series:
    dt = pd.to_datetime(s, errors="coerce", dayfirst=True, infer_datetime_format=True)
    if dt.notna().sum() < max(3, int(0.2 * len(s))):
        as_num = pd.to_numeric(s, errors="coerce")
        if as_num.notna().sum() > max(3, int(0.2 * len(s))):
            dt = excel_serial_to_datetime(s)
    return dt

def parse_dates(df: pd.DataFrame) -> pd.DataFrame:
    for c in DATE_COLS:
        if c in df.columns:
            df[c] = robust_parse_datetime(df[c])
    return df

def to_excel_download(df: pd.DataFrame) -> bytes:
    out = BytesIO()
    with pd.ExcelWriter(out, engine="xlsxwriter") as w:
        df.to_excel(w, index=False, sheet_name="datos_filtrados")
    return out.getvalue()

def safe_unique_opts(df, col):
    return sorted([x for x in df[col].dropna().astype(str).unique().tolist() if x != ""]) if col in df.columns else []

def agg_both(dff: pd.DataFrame, group_col: str, top_n: int | None = None) -> pd.DataFrame:
    """Long: [group_col, 'Métrica', 'Valor'] con Pacientes únicos vs Historias."""
    if group_col is None or group_col not in dff.columns:
        return pd.DataFrame(columns=[group_col or "Grupo","Métrica","Valor"])
    u = dff.groupby(group_col, dropna=False)["NUMERO_PACIENTE"].nunique(dropna=True).reset_index(name="Pacientes únicos")
    h = dff.dropna(subset=["NUMERO_PACIENTE"]).groupby(group_col, dropna=False)["NUMERO_PACIENTE"].size().reset_index(name="Historias")
    out = pd.merge(u, h, on=group_col, how="outer")
    for c in ["Pacientes únicos","Historias"]:
        if c in out.columns:
            out[c] = pd.to_numeric(out[c], errors="coerce").fillna(0)
    out[group_col] = out[group_col].fillna("Sin dato")
    out = out.sort_values("Historias", ascending=False)
    if top_n:
        out = out.head(top_n)
    return out.melt(id_vars=[group_col], value_vars=["Pacientes únicos","Historias"],
                    var_name="Métrica", value_name="Valor")

def build_age_buckets_both(dff: pd.DataFrame) -> pd.DataFrame:
    if "EDAD" not in dff.columns:
        return pd.DataFrame(columns=["RANGO_EDAD","Métrica","Valor"])
    edades = pd.to_numeric(dff["EDAD"], errors="coerce")
    bins = [-0.1, 5, 11, 17, 29, 44, 59, 74, 150]
    labels = ["0-5","6-11","12-17","18-29","30-44","45-59","60-74","75+"]
    tmp = dff.assign(EDAD_NUM=edades, RANGO_EDAD=pd.cut(edades, bins=bins, labels=labels))
    u = tmp.groupby("RANGO_EDAD")["NUMERO_PACIENTE"].nunique(dropna=True).reset_index(name="Pacientes únicos")
    h = tmp.dropna(subset=["NUMERO_PACIENTE"]).groupby("RANGO_EDAD")["NUMERO_PACIENTE"].size().reset_index(name="Historias")
    out = pd.merge(u, h, on="RANGO_EDAD", how="outer")
    for c in ["Pacientes únicos","Historias"]:
        if c in out.columns:
            out[c] = pd.to_numeric(out[c], errors="coerce").fillna(0)
    out["RANGO_EDAD"] = pd.Categorical(out["RANGO_EDAD"], categories=labels, ordered=True)
    out = out.sort_values("RANGO_EDAD")
    return out.melt(id_vars=["RANGO_EDAD"], value_vars=["Pacientes únicos","Historias"],
                    var_name="Métrica", value_name="Valor")

def total_by(dff: pd.DataFrame, cat: str) -> pd.DataFrame:
    if cat not in dff.columns:
        return pd.DataFrame(columns=[cat,"Valor"])
    return dff.dropna(subset=["NUMERO_PACIENTE"]).groupby(cat)["NUMERO_PACIENTE"].size().reset_index(name="Valor")


# ===================== CHARTS (web) =====================
def bar_two_metrics(df_long: pd.DataFrame, y_field: str, y_title: str):
    if df_long is None or df_long.empty:
        return None
    base = alt.Chart(df_long).encode(
        y=alt.Y(f"{y_field}:N", sort='-x', title=y_title,
                axis=alt.Axis(labelColor="#E5E7EB", titleColor="#F8FAFC")),
        x=alt.X("Valor:Q", title="Valor",
                axis=alt.Axis(labelColor="#E5E7EB", titleColor="#F8FAFC")),
        color=alt.Color("Métrica:N", scale=alt.Scale(range=[PALETTE[0], PALETTE[2]]),
                        legend=alt.Legend(title="Métrica", labelColor="#F8FAFC", titleColor="#F8FAFC")),
        tooltip=[f"{y_field}:N","Métrica:N",alt.Tooltip("Valor:Q", format=",.0f")]
    )
    bars = base.mark_bar()
    labels = base.mark_text(align="left", dx=6, color="#F8FAFC", fontSize=14, fontWeight="bold")\
                 .encode(text=alt.Text("Valor:Q", format=",.0f"))
    return (bars + labels).properties(height=380, background="#0f172a")

def donut_chart(df_counts: pd.DataFrame, cat_col: str, title: str, max_slices=12):
    if df_counts is None or df_counts.empty:
        return None
    dft = df_counts.sort_values("Valor", ascending=False).head(max_slices).copy()
    total = float(dft["Valor"].sum()) if len(dft) else 0.0
    dft["pct"] = np.where(total > 0, dft["Valor"]/total, 0)
    base = alt.Chart(dft).encode(
        theta=alt.Theta("Valor:Q", stack=True),
        color=alt.Color(f"{cat_col}:N", legend=alt.Legend(title=cat_col, labelColor="#F8FAFC", titleColor="#F8FAFC"),
                        scale=alt.Scale(range=PALETTE)),
        tooltip=[f"{cat_col}:N", alt.Tooltip("Valor:Q", format=",.0f"), alt.Tooltip("pct:Q", format=".1%")]
    )
    arc = base.mark_arc(innerRadius=60, outerRadius=120)
    text = base.mark_text(radius=140, color="#F8FAFC", fontSize=12, fontWeight="bold").encode(text=alt.Text(f"{cat_col}:N"))
    values = base.mark_text(radius=100, color="#F8FAFC", fontSize=12).encode(text=alt.Text("Valor:Q", format=",.0f"))
    return (arc + text + values).properties(height=360, title=title, background="#0f172a")


# ===================== CARGA DE DATOS =====================
@st.cache_data(show_spinner=False)
def load_data(uploaded_file: BytesIO) -> pd.DataFrame:
    name = uploaded_file.name.lower()
    if name.endswith(".csv"):
        df = pd.read_csv(uploaded_file, dtype=str, sep=None, engine="python")
    else:
        df = pd.read_excel(uploaded_file, dtype=str)
    df = normalize_cols(df)
    df = ensure_columns(df)
    df["EDAD"] = pd.to_numeric(df["EDAD"], errors="coerce")
    df = parse_dates(df)
    return df

# ===== NUEVO: lectura por URL (xlsx/xls/csv/parquet) =====
def _to_github_raw(url: str) -> str:
    """Convierte https://github.com/.../blob/... -> https://raw.githubusercontent.com/..."""
    if "github.com" in url and "/blob/" in url:
        owner_repo, path = url.split("github.com/")[1].split("/blob/")
        return f"https://raw.githubusercontent.com/{owner_repo}/{path}"
    return url

@st.cache_data(ttl=600, show_spinner=False)
def _fetch_bytes(url: str, token: str | None = None) -> bytes:
    headers = {"User-Agent": "streamlit-app"}
    if token:
        headers["Authorization"] = f"token {token}"
    r = requests.get(url, headers=headers, timeout=60)
    r.raise_for_status()
    return r.content

@st.cache_data(ttl=600, show_spinner=False)
def load_data_from_url(url: str, token: str | None = None) -> pd.DataFrame:
    url = _to_github_raw(url)
    raw = _fetch_bytes(url, token)
    bio = BytesIO(raw)
    lower = url.lower()

    # Detecta por extensión con fallback por contenido
    if lower.endswith(".parquet"):
        try:
            df = pd.read_parquet(bio, engine="pyarrow")
        except Exception:
            bio.seek(0)
            df = pd.read_csv(bio, dtype=str)
    elif lower.endswith(".csv"):
        try:
            df = pd.read_csv(bio, dtype=str)
        except Exception:
            bio.seek(0)
            df = pd.read_parquet(bio, engine="pyarrow")
    else:  # .xlsx / .xls (o sin extensión conocida)
        try:
            df = pd.read_excel(bio, dtype=str)
        except Exception:
            bio.seek(0)
            try:
                df = pd.read_csv(bio, dtype=str)
            except Exception:
                bio.seek(0)
                df = pd.read_parquet(bio, engine="pyarrow")

    df = normalize_cols(df)
    df = ensure_columns(df)
    df["EDAD"] = pd.to_numeric(df["EDAD"], errors="coerce")
    df = parse_dates(df)
    return df


# ===================== UI – SIDEBAR (FILTROS & AJUSTES) =====================
st.sidebar.title("📎 Carga de datos")

# ===== NUEVO: selector de origen =====
origen = st.sidebar.radio(
    "Origen de datos",
    ["Subir archivo", "URL (.xlsx/.xls/.csv/.parquet)"],
    index=0
)

df = None
if origen == "Subir archivo":
    uploaded = st.sidebar.file_uploader("Sube tu archivo (CSV/Excel)", type=["csv","xlsx","xls"])
    if uploaded is None:
        st.info("Sube un archivo o cambia a 'URL' para pegar un enlace (GitHub RAW, Drive público, etc.).")
        st.stop()
    df = load_data(uploaded)
else:
    url = st.sidebar.text_input(
        "Pega la URL del archivo",
        placeholder="https://raw.githubusercontent.com/usuario/repo/main/data/archivo.xlsx"
    )
    # Si tu repo es privado, define en Secrets: GITHUB_TOKEN="tu_PAT"
    gh_token = st.secrets.get("GITHUB_TOKEN", None)
    if not url:
        st.info("Pega una URL válida (RAW). Acepta .xlsx/.xls/.csv/.parquet.")
        st.stop()
    with st.spinner("Cargando datos desde URL..."):
        try:
            df = load_data_from_url(url, gh_token)
        except Exception as e:
            st.error(f"No se pudo leer la URL.\n\n{e}")
            st.stop()

with st.sidebar.expander("🎯 Segmentaciones", expanded=True):
    eps_f    = st.multiselect("EPS",          safe_unique_opts(df, "EPS"),           default=safe_unique_opts(df, "EPS"))
    esp_f    = st.multiselect("Especialidad", safe_unique_opts(df, "ESPECIALIDAD"),  default=safe_unique_opts(df, "ESPECIALIDAD"))
    prog_f   = st.multiselect("Programa",     safe_unique_opts(df, "PROGRAMA"),      default=safe_unique_opts(df, "PROGRAMA"))
    sede_f   = st.multiselect("Sede",         safe_unique_opts(df, "SEDE"),          default=safe_unique_opts(df, "SEDE"))
    genero_f = st.multiselect("Género",       safe_unique_opts(df, "GENERO"),        default=safe_unique_opts(df, "GENERO"))

    st.markdown("**Profesional**")
    available_prof_cols = [c for c in ["PRESTADOR","ENFERMERA"] if c in df.columns]
    if len(available_prof_cols)==0:
        chosen_prof_col, prof_f = None, []
        st.warning("No se encontró columna PRESTADOR ni ENFERMERA.")
    else:
        chosen_prof_col = st.selectbox("Columna", options=available_prof_cols, index=0)
        prof_f = st.multiselect(f"Selecciona {chosen_prof_col}",
                                safe_unique_opts(df, chosen_prof_col),
                                default=safe_unique_opts(df, chosen_prof_col))

    # Rango de fechas
    min_date = pd.to_datetime("1900-01-01")
    max_date = pd.to_datetime("today").normalize()
    if df["FECHA_ATENCION"].notna().any():
        min_date = df["FECHA_ATENCION"].dropna().min()
        max_date = df["FECHA_ATENCION"].dropna().max()
    date_range = st.date_input("Rango de FECHA_ATENCION",
                               value=(min_date.date() if pd.notna(min_date) else None,
                                      max_date.date() if pd.notna(max_date) else None))

with st.sidebar.expander("⚙️ Ajustes de visualización", expanded=False):
    top_n = st.slider("Top N para tortas (EPS/Programa/Sede)", 5, 20, 10, 1)
    ppt_bars_vertical = st.toggle("Barras VERTICALES en la PPT", value=False)

# Aplicar filtros
filt = pd.Series(True, index=df.index)
def safe_in(col, values):
    return df[col].isin(values) if (col in df.columns and len(values)>0) else pd.Series(True, index=df.index)
filt = (
    safe_in("EPS", eps_f) &
    safe_in("ESPECIALIDAD", esp_f) &
    safe_in("PROGRAMA", prog_f) &
    safe_in("SEDE", sede_f) &
    safe_in("GENERO", genero_f)
)
if chosen_prof_col is not None:
    filt = filt & safe_in(chosen_prof_col, prof_f)
if (
    isinstance(date_range, tuple) and len(date_range)==2 and date_range[0] and date_range[1]
    and df["FECHA_ATENCION"].notna().any()
):
    start, end = pd.to_datetime(date_range[0]), pd.to_datetime(date_range[1])
    filt = filt & (df["FECHA_ATENCION"].between(start, end))
dff = df[filt].copy()


# ===================== UI – CONTENIDO (TABS) =====================
st.markdown("<h1>🩺 Dashboard de Pacientes</h1>", unsafe_allow_html=True)

tab1, tab2, tab3, tab4, tab5 = st.tabs(["📌 Resumen", "📊 Barras", "🥧 Donas", "📋 Tabla", "📤 Exportar"])

with tab1:
    st.markdown("<div class='section-title'>Indicadores Clave</div>", unsafe_allow_html=True)
    total_unicos    = dff["NUMERO_PACIENTE"].nunique(dropna=True)
    total_historias = dff["NUMERO_PACIENTE"].notna().count()
    total_registros = len(dff)

    c1, c2, c3 = st.columns(3)
    with c1:
        st.markdown(f"<div class='kpi-card'><div class='kpi-title'>Pacientes únicos</div><div class='kpi-value'>{total_unicos:,}</div></div>", unsafe_allow_html=True)
    with c2:
        st.markdown(f"<div class='kpi-card'><div class='kpi-title'>Historias (no distintas)</div><div class='kpi-value'>{total_historias:,}</div></div>", unsafe_allow_html=True)
    with c3:
        st.markdown(f"<div class='kpi-card'><div class='kpi-title'>Registros filtrados</div><div class='kpi-value'>{total_registros:,}</div></div>", unsafe_allow_html=True)

    if total_registros == 0:
        st.warning("No hay registros con los filtros seleccionados.")

with tab2:
    st.markdown("<div class='section-title'>Barras comparativas (Pacientes únicos vs Historias)</div>", unsafe_allow_html=True)
    colA, colB = st.columns(2)
    with colA:
        g_genero = agg_both(dff, "GENERO")
        ch = bar_two_metrics(g_genero, "GENERO", "Género")
        if ch is not None: st.altair_chart(ch, use_container_width=True)
        else: st.info("Sin datos de género.")
    with colB:
        g_edad = build_age_buckets_both(dff)
        if not g_edad.empty:
            base = alt.Chart(g_edad).encode(
                x=alt.X("Valor:Q", title="Valor", axis=alt.Axis(labelColor="#E5E7EB", titleColor="#F8FAFC")),
                y=alt.Y("RANGO_EDAD:N", title="Rango de edad",
                        sort=list(g_edad["RANGO_EDAD"].astype(str).unique()),
                        axis=alt.Axis(labelColor="#E5E7EB", titleColor="#F8FAFC")),
                color=alt.Color("Métrica:N", scale=alt.Scale(range=[PALETTE[0], PALETTE[2]]),
                                legend=alt.Legend(title="Métrica", labelColor="#F8FAFC", titleColor="#F8FAFC")),
                tooltip=["RANGO_EDAD:N","Métrica:N",alt.Tooltip("Valor:Q", format=",.0f")]
            )
            bars = base.mark_bar()
            labels = base.mark_text(align="left", dx=6, color="#F8FAFC", fontSize=14, fontWeight="bold")\
                         .encode(text=alt.Text("Valor:Q", format=",.0f"))
            st.altair_chart((bars + labels).properties(height=380, background="#0f172a"), use_container_width=True)
        else:
            st.info("Sin datos de edad.")

    colC, colD = st.columns(2)
    with colC:
        g_sede = agg_both(dff, "SEDE", top_n=20)
        ch = bar_two_metrics(g_sede, "SEDE", "Sede")
        if ch is not None: st.altair_chart(ch, use_container_width=True)
        else: st.info("Sin datos de sede.")
    with colD:
        g_prof = agg_both(dff, "PRESTADOR" if "PRESTADOR" in dff.columns else (chosen_prof_col or "ENFERMERA"), top_n=20)
        ch = bar_two_metrics(g_prof, (chosen_prof_col or "Profesional"), "Profesional")
        if ch is not None: st.altair_chart(ch, use_container_width=True)
        else: st.info("Sin datos de profesional.")

with tab3:
    st.markdown("<div class='section-title'>Donas (Historias)</div>", unsafe_allow_html=True)
    dA, dB, dC = st.columns(3)
    with dA:
        df_gen = total_by(dff, "GENERO"); ch = donut_chart(df_gen, "GENERO", "Género (Historias)", max_slices=10)
        if ch is not None: st.altair_chart(ch, use_container_width=True)
        else: st.info("Sin datos de género.")
    with dB:
        df_eps = total_by(dff, "EPS"); ch = donut_chart(df_eps, "EPS", f"EPS Top {top_n} (Historias)", max_slices=top_n)
        if ch is not None: st.altair_chart(ch, use_container_width=True)
        else: st.info("Sin datos de EPS.")
    with dC:
        df_prog = total_by(dff, "PROGRAMA"); ch = donut_chart(df_prog, "PROGRAMA", f"Programa Top {top_n} (Historias)", max_slices=top_n)
        if ch is not None: st.altair_chart(ch, use_container_width=True)
        else: st.info("Sin datos de Programa.")

    dD, dE, _ = st.columns(3)
    with dD:
        df_sede = total_by(dff, "SEDE"); ch = donut_chart(df_sede, "SEDE", f"Sede Top {top_n} (Historias)", max_slices=top_n)
        if ch is not None: st.altair_chart(ch, use_container_width=True)
        else: st.info("Sin datos de Sede.")
    with dE:
        if chosen_prof_col:
            df_prof = total_by(dff, chosen_prof_col); ch = donut_chart(df_prof, chosen_prof_col, f"Profesional Top {top_n} (Historias)", max_slices=top_n)
            if ch is not None: st.altair_chart(ch, use_container_width=True)
            else: st.info("Sin datos de Profesional.")

with tab4:
    st.markdown("<div class='section-title'>Detalle (registros filtrados)</div>", unsafe_allow_html=True)
    st.dataframe(dff, use_container_width=True, hide_index=True)

with tab5:
    st.markdown("<div class='section-title'>Exportar archivos</div>", unsafe_allow_html=True)
    cexp1, cexp2 = st.columns(2)
    with cexp1:
        st.download_button("⬇️ Descargar CSV",
            data=dff.to_csv(index=False).encode("utf-8-sig"),
            file_name="filtrados.csv",
            mime="text/csv"
        )
        st.download_button("⬇️ Descargar Excel",
            data=to_excel_download(dff),
            file_name="filtrados.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

    # ---------------- PPTX ----------------
    def ppt_add_title(prs, title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def ppt_add_kpis(prs, unicos, historias, registros):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
        slide.shapes.title.text = "Indicadores Clave"
        left, top = Inches(0.7), Inches(1.6)
        tx = slide.shapes.add_textbox(left, top, Inches(9), Inches(3))
        tf = tx.text_frame; tf.clear()
        for label, val, color in [
            ("Pacientes únicos", unicos, RGBColor(96,165,250)),
            ("Historias (no distintas)", historias, RGBColor(52,211,153)),
            ("Registros filtrados", registros, RGBColor(251,191,36)),
        ]:
            p = tf.add_paragraph(); p.text = f"{label}: {val:,}"
            p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = color

    def ppt_add_bar(prs, title, df_long, cat_col, vertical=False):
        if df_long is None or df_long.empty: return
        pt = df_long.pivot_table(index=cat_col, columns="Métrica", values="Valor", aggfunc="sum").fillna(0)
        categories = pt.index.astype(str).tolist()
        data = CategoryChartData(); data.categories = categories
        if "Pacientes únicos" in pt.columns: data.add_series("Pacientes únicos", pt["Pacientes únicos"].tolist())
        if "Historias" in pt.columns: data.add_series("Historias", pt["Historias"].tolist())
        slide = prs.slides.add_slide(prs.slide_layouts[5]); slide.shapes.title.text = title
        x, y, cx, cy = Inches(0.7), Inches(1.7), Inches(9.0), Inches(4.3)
        chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED if vertical else XL_CHART_TYPE.BAR_CLUSTERED
        chart = slide.shapes.add_chart(chart_type, x, y, cx, cy, data).chart
        chart.has_title = False; chart.has_legend = True; chart.legend.include_in_layout = False
        chart.value_axis.has_major_gridlines = True
        chart.category_axis.tick_labels.font.size = Pt(10)
        chart.value_axis.tick_labels.font.size = Pt(10)

    def ppt_add_pie(prs, title, df_counts, cat_col, top=10):
        if df_counts is None or df_counts.empty: return
        dft = df_counts.sort_values("Valor", ascending=False).head(top)
        data = CategoryChartData()
        data.categories = dft[cat_col].astype(str).tolist()
        data.add_series("Historias", dft["Valor"].tolist())
        slide = prs.slides.add_slide(prs.slide_layouts[5]); slide.shapes.title.text = title
        x, y, cx, cy = Inches(0.7), Inches(1.7), Inches(9.0), Inches(4.3)
        chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, data).chart
        chart.has_legend = True; chart.legend.include_in_layout = False
        chart.plots[0].has_data_labels = True
        dl = chart.plots[0].data_labels; dl.number_format = '0'; dl.show_percentage = False; dl.font.size = Pt(10)

    def ppt_add_table(prs, title, df_sample):
        slide = prs.slides.add_slide(prs.slide_layouts[5]); slide.shapes.title.text = title
        rows, cols = df_sample.shape[0]+1, df_sample.shape[1]
        x, y, cx, cy = Inches(0.4), Inches(1.5), Inches(9.6), Inches(4.8)
        table = slide.shapes.add_table(rows, cols, x, y, cx, cy).table
        for j, col in enumerate(df_sample.columns):
            cell = table.cell(0, j); cell.text = str(col)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(10)
        for i in range(df_sample.shape[0]):
            for j in range(cols):
                cell = table.cell(i+1, j)
                val = df_sample.iat[i, j]
                cell.text = "" if pd.isna(val) else str(val)
                cell.text_frame.paragraphs[0].font.size = Pt(9)

    def build_ppt(dff: pd.DataFrame, top_n: int, chosen_prof_col: str | None, vertical_bars: bool) -> bytes:
        prs = Presentation()
        # Portada
        ppt_add_title(prs, "Dashboard de Pacientes — Informe", datetime.now().strftime("%d/%m/%Y %H:%M"))
        # KPIs
        unicos = dff["NUMERO_PACIENTE"].nunique(dropna=True)
        historias = dff["NUMERO_PACIENTE"].notna().count()
        registros = len(dff)
        ppt_add_kpis(prs, unicos, historias, registros)
        # Barras comparativas
        ppt_add_bar(prs, "Género — Pacientes únicos vs Historias", agg_both(dff, "GENERO"), "GENERO", vertical=vertical_bars)
        age_long = build_age_buckets_both(dff).rename(columns={"RANGO_EDAD":"RANGO EDAD"})
        ppt_add_bar(prs, "Rango de edad — Pacientes únicos vs Historias", age_long, "RANGO EDAD", vertical=vertical_bars)
        ppt_add_bar(prs, "Sede — Pacientes únicos vs Historias (Top 20)", agg_both(dff, "SEDE", top_n=20), "SEDE", vertical=vertical_bars)
        ppt_add_bar(prs, f"EPS — Pacientes únicos vs Historias (Top {top_n})", agg_both(dff, "EPS", top_n=top_n), "EPS", vertical=vertical_bars)
        ppt_add_bar(prs, f"Programa — Pacientes únicos vs Historias (Top {top_n})", agg_both(dff, "PROGRAMA", top_n=top_n), "PROGRAMA", vertical=vertical_bars)
        if chosen_prof_col:
            ppt_add_bar(prs, f"Profesional — Pacientes únicos vs Historias (Top 20) [{chosen_prof_col}]",
                        agg_both(dff, chosen_prof_col, top_n=20), chosen_prof_col, vertical=vertical_bars)
        # Donas (Historias)
        ppt_add_pie(prs, "Género — Historias", total_by(dff, "GENERO"), "GENERO", top=10)
        ppt_add_pie(prs, f"EPS — Historias (Top {top_n})", total_by(dff, "EPS"), "EPS", top=top_n)
        ppt_add_pie(prs, f"Programa — Historias (Top {top_n})", total_by(dff, "PROGRAMA"), "PROGRAMA", top=top_n)
        ppt_add_pie(prs, f"Sede — Historias (Top {top_n})", total_by(dff, "SEDE"), "SEDE", top=top_n)
        # Muestra de tabla
        sample_cols = [c for c in dff.columns if c in ["NUMERO_PACIENTE","NOMBRE_COMPLETO","GENERO","EDAD","SEDE","EPS","PROGRAMA","ESPECIALIDAD","FECHA_ATENCION"]]
        if not sample_cols:
            sample_cols = dff.columns.tolist()[:8]
        sample = dff[sample_cols].copy().head(15).fillna("")
        ppt_add_table(prs, "Muestra de registros filtrados", sample)
        # Export
        bio = BytesIO(); prs.save(bio); bio.seek(0); return bio.getvalue()

    ppt_bytes = build_ppt(dff, top_n, chosen_prof_col if 'chosen_prof_col' in locals() else None, vertical_bars=ppt_bars_vertical)
    st.download_button(
        "📥 Descargar presentación (PPTX)",
        data=ppt_bytes,
        file_name="Dashboard_Pacientes.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
